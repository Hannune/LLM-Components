import fitz  # PyMuPDF
from docx import Document
import pptx, docx
import re
import pdfplumber


class naive_text_extraction:
    def __init__(self, logger):
        self.file_instance = None
        self.extracted_text = {}
        self.line_num_threshold = 3
        self.space_num_threshold = 3
        self.logger = logger


    def start(self, file_instance, doc_title, doc_type):
        
        self.file_instance = file_instance
        self.extracted_text = {}
        self.extracted_text["contents"] = {}
        self.extracted_text["document_title"] = doc_title
        self.extracted_text["document_type"] = doc_type

        if doc_type.lower() == "pdf":
            self.parse_pdf()
        elif doc_type.lower() == "docx":
            self.parse_docx()
        elif doc_type.lower() in ["pptx", "ppt"]:
            self.parse_ppt()   

        extracted_text = self.extracted_text
        self.extracted_text = {}
        self.file_instance = None

        print(f"Finished text extraction")
        return extracted_text
        

    def parse_pdf(self):
        # doc = fitz.open(stream=self.file_instance, filetype="pdf")
        # self.extracted_text["contents_index_type"] = "page_number"

        # # option 1
        # for page in doc:
        #     page_num = doc.index(page) + 1
        #     text_dict = page.get_text("dict")  # Text as dictionary
        #     page_text = ""
        #     for block in text_dict['blocks']:
        #         try:
        #             for line in block['lines']:
        #                 for span in line['spans']:
        #                     page_text += span['text'] + "\n"
        #             page_text += "\n"
        #         except:
        #             continue
        #     page_text = self.remove_excessive_lines_and_spaces(page_text)
        #     self.extracted_text["contents"][page_num] = {"text": page_text.strip()}

        # # option 2
        # for page in doc:
        #     page_num = page.number + 1
        #     blocks = page.get_text("blocks")  # Extract text in blocks
        #     page_text = ""
        #     for block in blocks:
        #         page_text += block[4] + " "  # block[4] contains the text in each block
        #     page_text = self.remove_excessive_lines_and_spaces(page_text)
        #     self.extracted_text["contents"][page_num] = {"text": page_text}
        
        self.extracted_text["contents_index_type"] = "page_number"
        with pdfplumber.open(self.file_instance) as pdf:
            for page in pdf.pages:
                page_num = page.page_number
                text_blocks = page.extract_text()
                if text_blocks:
                    page_text = self.remove_excessive_lines_and_spaces(text_blocks)
                else:
                    page_text = ""  # Handle pages that may be empty or have no text
                self.extracted_text["contents"][page_num] = {"text": page_text}


    def parse_docx(self):
        """
        Unlike pdf and pptx, docx does not have a built-in function to extract page number
        Also when extracting text from docx file, it uses paragraph as a unit
        But these paragraphs does not align with actual paragraphs in the document
        So, I will use whole text as a single paragraph, and store it on the first index of the contents
        """
        doc = Document(self.file_instance)

        self.extracted_text["contents_index_type"] = "whole_text_or_table_number"

        cnt_text_idx = 1
        cnt_table_idx = 1
        
        full_text = ""
        for para in doc.paragraphs:
            text = para.text
            if text in ["", " ", None]:
                continue

            # option 1
            full_text += text + "\n"
            # option 2
            # text = self.remove_excessive_lines_and_spaces(text)
            # if cnt_text_idx in self.extracted_text["contents"]:
            #     self.extracted_text["contents"][cnt_text_idx]["text"] = text
            # else:
            #     self.extracted_text["contents"][cnt_text_idx] = {"text": text}            
            # cnt_text_idx += 1
            
        # option 1
        full_text = self.remove_excessive_lines_and_spaces(full_text)
        self.extracted_text["contents"][cnt_text_idx] = {"text": full_text}

        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(",".join(row_data))
            csv_table = "\n".join(table_data)
            
            if cnt_table_idx in self.extracted_text["contents"]:
                self.extracted_text["contents"][cnt_table_idx]["table"] = csv_table
            else:
                self.extracted_text["contents"][cnt_table_idx] = {"table": csv_table}
            cnt_table_idx += 1


    def parse_ppt(self):
        presentation = pptx.Presentation(self.file_instance)

        self.extracted_text["contents_index_type"] = "slide_number"
        
        for slide in presentation.slides:
            # Get slide number of the current slide
            slide_number = presentation.slides.index(slide) + 1
            self.extracted_text["contents"][slide_number] = {}
            texts = ""
            tables = ""
            for shape in slide.shapes:
                # Handling text in both placeholders and other text-containing shapes
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = ""
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + "\n"
                    texts += text + "\n\n"

                
                # Check if the shape is a table and extract data
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text_frame.text if cell.text_frame else ""
                            row_data.append(cell_text)
                        table_data.append(",".join(row_data))
                    # Joining rows with newline to format the full table as CSV
                    csv_table = "\n".join(table_data)
                    tables += csv_table + "\n\n"
            if texts:
                texts = self.remove_excessive_lines_and_spaces(texts)
                self.extracted_text["contents"][slide_number]["text"] = texts
            if tables:
                self.extracted_text["contents"][slide_number]["table"] = tables


    def remove_excessive_lines_and_spaces(self, text):
        # Replace 3 or more consecutive newlines with 3 newlines
        text = re.sub(r'\n{'+ str(self.line_num_threshold) + ',}', '\n'*self.line_num_threshold, text)
        # Replace 3 or more consecutive spaces with 3 spaces
        text = re.sub(r' {'+ str(self.space_num_threshold) + ',}', ' '*self.space_num_threshold, text)

        return text

















